#!/usr/bin/env python3
"""
Convert PDF slides to PowerPoint presentation at 300 DPI.
"""

from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import os

def pdf_to_pptx(pdf_path, output_path, dpi=300):
    """
    Convert a PDF file to PPTX with specified DPI.

    Args:
        pdf_path: Path to input PDF file
        output_path: Path for output PPTX file
        dpi: Resolution for image conversion (default: 300)
    """
    print(f"Converting {pdf_path} to {output_path} at {dpi} DPI...")

    # Convert PDF pages to images at specified DPI
    print(f"Rendering PDF pages at {dpi} DPI...")
    images = convert_from_path(pdf_path, dpi=dpi)
    print(f"Rendered {len(images)} pages")

    # Create PowerPoint presentation
    prs = Presentation()

    # Set slide dimensions to match standard presentation (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add each image as a slide
    for i, image in enumerate(images):
        print(f"Adding slide {i+1}/{len(images)}")

        # Add blank slide
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Save image temporarily
        temp_image = f"temp_slide_{i}.png"
        image.save(temp_image, "PNG")

        # Add image to slide, filling the entire slide
        left = Inches(0)
        top = Inches(0)
        slide.shapes.add_picture(temp_image, left, top,
                                width=prs.slide_width,
                                height=prs.slide_height)

        # Remove temporary image
        os.remove(temp_image)

    # Save presentation
    prs.save(output_path)
    print(f"Successfully created {output_path}")

if __name__ == "__main__":
    pdf_path = "slides.pdf"
    output_path = "slides.pptx"
    pdf_to_pptx(pdf_path, output_path, dpi=300)
